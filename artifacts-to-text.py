#!/usr/bin/env python3
"""
artifacts-to-text.py

Converts mixed-format SDR artifacts into plain-text files and generates
a comprehensive conversion manifest.
"""

import os
import sys
import argparse
import json
import datetime
import traceback
from pathlib import Path

# Tracking available dependencies
DEPS = {
    "pdfplumber": False,
    "pdf2image": False,
    "python-docx": False,
    "python-pptx": False,
    "extract-msg": False,
    "easyocr": False,
    "pillow": False
}

# Dynamic Imports & Dependency Assertions
try:
    import pdfplumber
    DEPS["pdfplumber"] = True
except ImportError:
    pass

try:
    import pdf2image
    DEPS["pdf2image"] = True
except ImportError:
    pass

try:
    import docx
    DEPS["python-docx"] = True
except ImportError:
    pass

try:
    import pptx
    DEPS["python-pptx"] = True
except ImportError:
    pass

try:
    import extract_msg
    DEPS["extract-msg"] = True
except ImportError:
    pass

try:
    import easyocr
    DEPS["easyocr"] = True
except ImportError:
    pass

try:
    import PIL
    from PIL import Image
    DEPS["pillow"] = True
except ImportError:
    pass


class ArtifactConverter:
    def __init__(self, quiet=False):
        self.quiet = quiet
        self.reader = None
        
        # Instantiate OCR engine only if dependency exists and is required
        if DEPS["easyocr"]:
            if not self.quiet:
                print("[*] Initializing EasyOCR Engine (English)...")
            try:
                self.reader = easyocr.Reader(['en'], gpu=False)
            except Exception as e:
                if not self.quiet:
                    print(f"[-] Failed to initialize EasyOCR: {e}")

    def log(self, message):
        if not self.quiet:
            print(message)

    def check_dependencies(self):
        print("Checking dependencies...")
        all_available = True
        for dep, available in DEPS.items():
            status = "✓" if available else "✗"
            print(f"  {status} {dep}")
            if not available:
                all_available = False
        
        if all_available:
            print("✓ All dependencies available")
            return True
        else:
            print("✗ Missing dependencies. Please run:")
            print("  pip install pdfplumber pdf2image python-docx python-pptx extract-msg pillow easyocr")
            return False

    def extract_pdf(self, file_path):
        text_content = []
        embedded_image_count = 0
        
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text_content.append(f"--- PAGE {i} ---")
                
                # Extract text labels / native text
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_content.append(page_text.strip())
                
                # Check for native tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        text_content.append("\n[TABLE]")
                        for row in table:
                            # Clean up None values and join items cleanly
                            row_str = " | ".join([str(cell).strip() if cell is not None else "" for cell in row])
                            text_content.append(row_str)
                        text_content.append("[END TABLE]\n")
                
                # Limitations note: Record metadata presence of images rather than exporting[cite: 2]
                if hasattr(page, 'images') and page.images:
                    embedded_image_count += len(page.images)
                    text_content.append(f"\n[IMAGE FOUND ON PAGE {i} - See converted images]")

        return "\n".join(text_content), embedded_image_count

    def extract_docx(self, file_path):
        doc = docx.Document(file_path)
        text_content = []
        embedded_image_count = 0
        
        # Paragraph extraction
        for p in doc.paragraphs:
            if p.text.strip():
                text_content.append(p.text.strip())
                
        # Tables extraction
        for table in doc.tables:
            text_content.append("\n[TABLE]")
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                text_content.append(" | ".join(row_text))
            text_content.append("[END TABLE]\n")
            
        # Limitations note: Scan document structure for inline shapes/images[cite: 2]
        for rel in doc.part.relations.values():
            if "image" in rel.target_ref:
                embedded_image_count += 1
                
        if embedded_image_count > 0:
            text_content.append(f"\n[Embedded Images Encountered in Document Structural Layout: {embedded_image_count}]")
            
        return "\n".join(text_content), embedded_image_count

    def extract_pptx(self, file_path):
        prs = pptx.Presentation(file_path)
        text_content = []
        
        for i, slide in enumerate(prs.slides, start=1):
            text_content.append(f"--- Slide {i} ---")
            
            # Extract standard shape text and tables
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_content.append(paragraph.text.strip())
                if shape.has_table:
                    text_content.append("\n[TABLE]")
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        text_content.append(" | ".join(row_text))
                    text_content.append("[END TABLE]\n")
            
            # Extract Speaker Notes[cite: 2]
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_content.append(f"\n[SPEAKER NOTES]\n{notes}\n")
                    
        return "\n".join(text_content), 0

    def extract_msg(self, file_path):
        msg = extract_msg.Message(file_path)
        text_content = [
            f"Subject: {msg.subject}",
            f"From: {msg.sender}",
            f"To: {msg.to}",
            f"Date: {msg.date}",
            "\n--- Body ---",
            msg.body if msg.body else "[No Text Body Available]"
        ]
        
        # Track and note attachments[cite: 2]
        if msg.attachments:
            text_content.append("\n--- Attachments Checklist ---")
            for attachment in msg.attachments:
                filename = attachment.longFilename or attachment.shortFilename or "Unknown_Attachment"
                text_content.append(f"[Attached File]: {filename}")
                
        msg.close()
        return "\n".join(text_content), len(msg.attachments) if msg.attachments else 0

    def extract_ocr(self, file_path):
        if not self.reader:
            return "[WARNING: EasyOCR framework skipped or unavailable for processing]", 0
        try:
            results = self.reader.readtext(str(file_path), detail=0)
            if results:
                return "\n".join(results), 0
            return "[No OCR text labels detected in source layout matrix]", 0
        except Exception as e:
            return f"[OCR processing encountered error structural evaluation bounds: {e}]", 0

    def process_artifacts(self, input_dir, output_dir):
        input_path = Path(input_dir)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        manifest = {
            "conversion_date": datetime.datetime.now().isoformat(),
            "files_processed": [],
            "errors": [],
            "coverage": {}
        }
        
        # Supported Extension Routing Matrix
        supported_exts = {
            ".pdf": ("pdfplumber", self.extract_pdf),
            ".docx": ("python-docx", self.extract_docx),
            ".pptx": ("python-pptx", self.extract_pptx),
            ".msg": ("extract-msg", self.extract_msg),
            ".png": ("easyocr", self.extract_ocr),
            ".jpg": ("easyocr", self.extract_ocr),
            ".jpeg": ("easyocr", self.extract_ocr),
            ".txt": (None, lambda p: (p.read_text(encoding="utf-8", errors="replace"), 0))
        }
        
        # Verify mandatory tools for present files before processing[cite: 2]
        found_extensions = set(p.suffix.lower() for p in input_path.rglob("*") if p.is_file())
        for ext in found_extensions:
            if ext in supported_exts:
                required_dep = supported_exts[ext][0]
                if required_dep and not DEPS[required_dep]:
                    print(f"ERROR: Required dependency '{required_dep}' for configuration extension {ext} missing.")
                    sys.exit(1)

        # Scans directory recursively[cite: 2]
        all_files = [p for p in input_path.rglob("*") if p.is_file()]
        supported_files = [p for p in all_files if p.suffix.lower() in supported_exts]
        
        if not supported_files:
            print("[-] Error processing parameters: No supported artifact files found in the specified target directory.")
            sys.exit(1)
            
        self.log(f"[+] Found {len(supported_files)} target files. Beginning extraction sequence...")
        
        for file_p in supported_files:
            ext = file_p.suffix.lower()
            self.log(f"[*] Processing: {file_p.name} ({ext})")
            
            # Setup output filename logic matching original stem pattern bounds[cite: 1, 2]
            out_filename = f"{file_p.stem}_extracted.txt"
            target_out_path = output_path / out_filename
            
            _, handler_func = supported_exts[ext]
            
            try:
                extracted_text, _ = handler_func(file_p)
                
                # Write UTF-8 artifact text output[cite: 2]
                target_out_path.write_text(extracted_text, encoding="utf-8")
                
                # Update manifest statistics metrics[cite: 2]
                manifest["files_processed"].append({
                    "original": str(file_p.relative_to(input_path.parent) if input_path.parent != input_path else file_p),
                    "output": str(target_out_path),
                    "format": ext
                })
                manifest["coverage"][ext] = manifest["coverage"].get(ext, 0) + 1
                
            except Exception as e:
                err_msg = f"Extraction failure on {file_p.name}: {str(e)}"
                self.log(f"[-] {err_msg}")
                manifest["errors"].append({
                    "file": str(file_p),
                    "error": err_msg,
                    "traceback": traceback.format_exc()
                })
                
        # Write CONVERSION_MANIFEST.json summary[cite: 1, 2]
        manifest_path = output_path / "CONVERSION_MANIFEST.json"
        manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
        
        print(f"[+] Operational layout finalized. Unified log exported safely: {output_path}")
        print(f"Successful processes tracked across coverage matrices: {json.dumps(manifest['coverage'])}")
        if manifest["errors"]:
            print(f"[!] Encountered {len(manifest['errors'])} conversion tracking exception anomalies. See manifest for logs.")


def main():
    parser = argparse.ArgumentParser(description="Secure Design Review (SDR) Artifact-to-Text Utility Engine")
    parser.add_argument("artifacts_dir", nargs="?", help="Input folder containing target artifacts to convert")
    parser.add_argument("--output", default="artifacts_converted", help="Destination folder tracking extraction outputs (default: artifacts_converted)")
    parser.add_argument("--check-deps", action="store_true", help="Checks dependency availability status and exits cleanly")
    parser.add_argument("--quiet", action="store_true", help="Suppresses verbose per-file terminal metrics logger flags")
    
    args = parser.parse_args()
    
    converter = ArtifactConverter(quiet=args.quiet)
    
    if args.check_deps:
        success = converter.check_dependencies()
        sys.exit(0 if success else 1)
        
    if not args.artifacts_dir:
        print("ERROR: Missing required positional directory argument: 'artifacts_dir'")
        parser.print_help()
        sys.exit(1)
        
    input_path = Path(args.artifacts_dir)
    if not input_path.exists() or not input_path.is_dir():
        print(f"[-] Fault: Input target directory structure '{args.artifacts_dir}' can't be reached or found.")
        sys.exit(1)
        
    converter.process_artifacts(args.artifacts_dir, args.output)


if __name__ == "__main__":
    main()
